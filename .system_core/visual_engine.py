"""
visual_engine.py — Nano Banana 2 Visual Production Engine (Global Core)
═══════════════════════════════════════════════════════════════════════════
.system_core | Antigravity V3.0 | Venture Studio Inheritance Engine

Consolidates document creation, brand-aware templating, and AI image
generation into a single engine. Inherits FileFactory patterns from
Alpha_V2_Genesis/skills but adds brand_identity.json enforcement.

Supports: PDF, DOCX, XLSX, PPTX, TXT, PNG (via DALL-E / local)
"""

import os
import sys
import json
import logging
from datetime import datetime

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
FACTORY_DIR = os.path.normpath(os.path.join(SCRIPT_DIR, ".."))
sys.path.insert(0, FACTORY_DIR)

logger = logging.getLogger("system_core.visual_engine")


class VisualEngine:
    """
    Nano Banana 2 — Brand-aware document & visual asset generator.

    Usage:
        from system_core import VisualEngine
        engine = VisualEngine(output_dir="./output", brand_file="soul/brand_identity.json")
        engine.create_pdf("Report content...", "quarterly_report.pdf")
        engine.generate_image("Hero banner for landing page", "hero.png")
    """

    def __init__(self, output_dir=None, brand_file=None):
        self.output_dir = output_dir or os.path.join(FACTORY_DIR, "generated_files")
        os.makedirs(self.output_dir, exist_ok=True)

        # Load brand identity if available
        self.brand = self._load_brand(brand_file)
        logger.info("VisualEngine initialized (brand: %s)", "loaded" if self.brand else "default")

    def _load_brand(self, brand_file):
        """Load brand_identity.json for style enforcement."""
        if brand_file and os.path.exists(brand_file):
            try:
                with open(brand_file, "r", encoding="utf-8") as f:
                    return json.load(f)
            except Exception as e:
                logger.warning("Failed to load brand file: %s", e)
        return None

    # ── Document Creation ────────────────────────────────

    def create_file(self, content, file_type, file_name):
        """
        Generate a file with brand enforcement.
        file_type: 'pdf', 'docx', 'xlsx', 'pptx', 'txt'
        """
        try:
            file_path = os.path.join(self.output_dir, file_name)
            ft = file_type.lower()

            if ft == "pdf":
                self._create_pdf(content, file_path)
            elif ft == "docx":
                self._create_docx(content, file_path)
            elif ft == "xlsx":
                self._create_xlsx(content, file_path)
            elif ft == "pptx":
                self._create_presentation(content, file_path)
            else:
                with open(file_path, "w", encoding="utf-8") as f:
                    f.write(str(content))

            logger.info("Created %s: %s", ft, file_name)
            return file_path
        except Exception as e:
            return f"File Generation Error: {str(e)}"

    def generate_image(self, prompt, file_name="image.png"):
        """Generate image via DALL-E 3, applying brand color context."""
        try:
            from openai import OpenAI
            import requests

            # Inject brand context into prompt
            brand_prompt = prompt
            if self.brand:
                colors = self.brand.get("colors", {})
                mission = self.brand.get("mission", "")
                if colors:
                    color_str = ", ".join(f"{k}: {v}" for k, v in colors.items())
                    brand_prompt = f"{prompt}. Use brand colors: {color_str}."
                if mission:
                    brand_prompt += f" Brand essence: {mission}"

            client = OpenAI()
            response = client.images.generate(
                model="dall-e-3",
                prompt=brand_prompt,
                size="1024x1024",
                quality="standard",
                n=1,
            )

            image_url = response.data[0].url
            file_path = os.path.join(self.output_dir, file_name)
            img_data = requests.get(image_url).content
            with open(file_path, "wb") as handler:
                handler.write(img_data)

            logger.info("Generated image: %s", file_name)
            return file_path
        except Exception as e:
            return f"Image Generation Error: {str(e)}"

    # ── Brand Metadata ────────────────────────────────────

    def get_brand_header(self):
        """Return brand header text for documents."""
        if not self.brand:
            return "Antigravity Venture Studio"
        return self.brand.get("company_name", self.brand.get("name", "Venture Studio"))

    def get_brand_colors(self):
        """Return brand color palette dict."""
        if not self.brand:
            return {"primary": "#3b82f6", "secondary": "#8b5cf6", "accent": "#06b6d4"}
        return self.brand.get("colors", {"primary": "#3b82f6", "secondary": "#8b5cf6"})

    def get_brand_fonts(self):
        """Return brand font families."""
        if not self.brand:
            return {"heading": "Inter", "body": "Inter", "mono": "JetBrains Mono"}
        return self.brand.get("fonts", {"heading": "Helvetica", "body": "Helvetica"})

    # ── Private Generators ────────────────────────────────

    def _create_pdf(self, content, path):
        """PDF with brand-aware header."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas as pdf_canvas

            c = pdf_canvas.Canvas(path, pagesize=letter)
            width, height = letter

            # Brand header
            fonts = self.get_brand_fonts()
            c.setFont("Helvetica-Bold", 14)
            c.drawString(40, height - 30, self.get_brand_header())
            c.setFont("Helvetica", 8)
            c.drawString(40, height - 44, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
            c.line(40, height - 50, width - 40, height - 50)

            # Body
            text_object = c.beginText(40, height - 70)
            text_object.setFont("Helvetica", 10)

            for line in content.split("\n"):
                if text_object.getY() < 40:
                    c.drawText(text_object)
                    c.showPage()
                    text_object = c.beginText(40, height - 40)
                    text_object.setFont("Helvetica", 10)
                text_object.textLine(line)

            c.drawText(text_object)
            c.save()
        except ImportError:
            # Fallback: plain text
            with open(path.replace(".pdf", ".txt"), "w", encoding="utf-8") as f:
                f.write(content)
            return path.replace(".pdf", ".txt")

    def _create_docx(self, content, path):
        """DOCX with brand heading."""
        try:
            from docx import Document
            doc = Document()
            doc.add_heading(self.get_brand_header(), level=0)
            for line in content.split("\n"):
                doc.add_paragraph(line)
            doc.save(path)
        except ImportError:
            with open(path.replace(".docx", ".txt"), "w", encoding="utf-8") as f:
                f.write(content)

    def _create_xlsx(self, content, path):
        """Excel from JSON data."""
        try:
            import pandas as pd
            data = []
            if isinstance(content, str):
                try:
                    data = json.loads(content)
                except Exception:
                    data = [{"Error": "Could not parse JSON for Excel"}]
            elif isinstance(content, list):
                data = content
            if not data:
                data = [{"Status": "No Data Provided"}]
            df = pd.DataFrame(data)
            df.to_excel(path, index=False)
        except ImportError:
            with open(path.replace(".xlsx", ".csv"), "w", encoding="utf-8") as f:
                f.write(str(content))

    def _create_presentation(self, content, path):
        """PPTX with brand title slide."""
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()
            try:
                data = json.loads(content) if isinstance(content, str) else content
            except Exception:
                data = {"title": "Error Parsing JSON", "slides": []}

            # Title Slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = data.get("title", self.get_brand_header())
            title_slide.placeholders[1].text = f"Generated by {self.get_brand_header()}"

            # Content Slides
            for slide_data in data.get("slides", []):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_data.get("title", "Slide")
                slide.placeholders[1].text_frame.text = slide_data.get("content", "")

                img_path = slide_data.get("image_path")
                if img_path and os.path.exists(img_path):
                    slide.shapes.add_picture(img_path, Inches(5.5), Inches(2), width=Inches(4))

            prs.save(path)
        except ImportError:
            with open(path.replace(".pptx", ".txt"), "w", encoding="utf-8") as f:
                f.write(str(content))


if __name__ == "__main__":
    engine = VisualEngine()
    print(f"VisualEngine initialized")
    print(f"  Output dir: {engine.output_dir}")
    print(f"  Brand: {engine.get_brand_header()}")
    print(f"  Colors: {engine.get_brand_colors()}")
    print(f"  Fonts: {engine.get_brand_fonts()}")

    # Quick test: create a text file
    result = engine.create_file("Hello from Nano Banana 2!", "txt", "test_visual.txt")
    print(f"  Test file: {result}")
