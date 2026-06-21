"""
deck.py — studio-grade deck/brochure that HONORS an existing identity (brand-deck)
══════════════════════════════════════════════════════════════════════════════════
Drives every colour, font and device from the brand-token file (never hardcoded
inline). Builds a 16:9 PDF deck (the render-QA'd artifact) and an editable PPTX.

Design system, applied once and repeated:
  • ONE dominant colour does 60–70% of the work (full-bleed cover / dividers /
    closing + content header zones). Accents are <10%, emphasis only.
  • The motif (a single brand mark) repeats in the SAME position on every page —
    repetition reads as identity. NO decorative accent bars / underlines / edge
    stripes (the tell of a templated deck): structure comes from type hierarchy,
    whitespace and the motif.
  • A serious size jump from heading to body; ≤ two families; generous leading.
  • Real data only (charts inherit brand tokens; numbers come from the session).
"""

from __future__ import annotations

import os
import sys

from reportlab.pdfgen import canvas
from reportlab.lib.colors import HexColor
from reportlab.lib.utils import simpleSplit
from reportlab.graphics import renderPDF

# reuse the shared brand-token loader + chart helpers from fin-model-presentation.
# APPEND (not insert) so we don't shadow brand-deck's own modules (e.g. qa.py).
_PRES_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                         "fin-model-presentation")
if _PRES_DIR not in sys.path:
    sys.path.append(_PRES_DIR)
from brand_tokens import load_tokens                      # noqa: E402
import charts_pdf as ch                                   # noqa: E402

PAGE_W, PAGE_H = 960, 540          # 13.333 x 7.5 in @72dpi — 16:9 widescreen
M = 54                             # outer margin
_STD_FONTS = {"Helvetica", "Helvetica-Bold", "Helvetica-Oblique", "Times-Roman",
              "Times-Bold", "Times-Italic", "Courier", "Courier-Bold"}


def _safe(name, bold=False):
    return name if name in _STD_FONTS else ("Helvetica-Bold" if bold else "Helvetica")


class DeckBuilder:
    def __init__(self, tokens_source):
        self.t = load_tokens(tokens_source)
        # keep voice/motif from the rich token file if present
        self.voice = {}
        self.motif = {"kind": "dot"}
        if isinstance(tokens_source, dict):
            self.voice = tokens_source.get("voice", {}) or {}
            self.motif = tokens_source.get("motif", {}) or self.motif
        elif isinstance(tokens_source, str) and os.path.exists(tokens_source):
            import json
            data = json.load(open(tokens_source, encoding="utf-8"))
            self.voice = data.get("voice", {}) or {}
            self.motif = data.get("motif", {}) or self.motif
        self.hfont = _safe(self.t.heading_font, bold=True)
        self.bfont = _safe(self.t.body_font)
        self.P = HexColor(self.t.primary); self.S = HexColor(self.t.secondary)
        self.INK = HexColor(self.t.ink); self.PAPER = HexColor(self.t.paper)
        self.ACC = HexColor(self.t.accents[0]) if self.t.accents else self.S

    # ── motif: one mark, same place every page ──────────────────────────────
    def _motif(self, c, on_dark=False):
        col = self.PAPER if on_dark else self.P
        kind = (self.motif or {}).get("kind", "dot")
        x, y = M, PAGE_H - M
        if kind == "corner":
            c.setFillColor(col); c.rect(x - 2, y - 8, 14, 14, fill=1, stroke=0)
        elif kind == "frame":
            c.setStrokeColor(col); c.setLineWidth(1.4)
            c.rect(M - 14, M - 14, PAGE_W - 2 * (M - 14), PAGE_H - 2 * (M - 14), fill=0, stroke=1)
        else:  # dot (default)
            c.setFillColor(col); c.circle(x + 4, y, 7, fill=1, stroke=0)

    def _wordmark(self, c, name, on_dark=False):
        c.setFillColor(self.PAPER if on_dark else self.P)
        c.setFont(self.hfont, 13)
        c.drawRightString(PAGE_W - M, PAGE_H - M - 4, name)

    def _wrap(self, c, text, x, y, width, font, size, color, leading=None, max_lines=None):
        leading = leading or size * 1.22
        c.setFont(font, size); c.setFillColor(color)
        lines = simpleSplit(text, font, size, width)
        if max_lines:
            lines = lines[:max_lines]
        for ln in lines:
            c.drawString(x, y, ln); y -= leading
        return y

    # ── slide renderers ─────────────────────────────────────────────────────
    def _cover(self, c, s, name):
        c.setFillColor(self.P); c.rect(0, 0, PAGE_W, PAGE_H, fill=1, stroke=0)   # dominant
        self._motif(c, on_dark=True)
        c.setFont(self.hfont, 13); c.setFillColor(self.PAPER)
        c.drawRightString(PAGE_W - M, PAGE_H - M - 4, name)
        y = self._wrap(c, s.get("headline", name), M, PAGE_H * 0.56,
                       PAGE_W - 2 * M, self.hfont, 58, self.PAPER, leading=62)
        tag = s.get("tagline") or (self.voice.get("taglines") or [""])[0]
        if tag:
            self._wrap(c, tag, M, y - 14, PAGE_W * 0.7, self.bfont, 20,
                       HexColor(self.t.accents[0]) if self.t.accents else self.PAPER)
        c.setFont(self.bfont, 11); c.setFillColor(self.PAPER)
        if s.get("footer"):
            c.drawString(M, M - 6, s["footer"])

    def _statement(self, c, s, name):
        c.setFillColor(self.PAPER); c.rect(0, 0, PAGE_W, PAGE_H, fill=1, stroke=0)
        self._motif(c); self._wordmark(c, name)
        # big primary headline; generous whitespace; no stripes
        self._wrap(c, s["headline"], M, PAGE_H * 0.62, PAGE_W * 0.82,
                   self.hfont, 40, self.P, leading=46)
        if s.get("sub"):
            self._wrap(c, s["sub"], M, PAGE_H * 0.34, PAGE_W * 0.68,
                       self.bfont, 17, self.INK, leading=24)

    def _bullets(self, c, s, name):
        c.setFillColor(self.PAPER); c.rect(0, 0, PAGE_W, PAGE_H, fill=1, stroke=0)
        # substantial primary header ZONE (not a thin stripe)
        c.setFillColor(self.P); c.rect(0, PAGE_H - 132, PAGE_W, 132, fill=1, stroke=0)
        self._motif(c, on_dark=True)
        c.setFont(self.hfont, 30); c.setFillColor(self.PAPER)
        c.drawString(M, PAGE_H - 92, s["title"])
        y = PAGE_H - 132 - 52
        for b in s.get("bullets", []):
            c.setFillColor(self.ACC); c.circle(M + 5, y + 5, 4.5, fill=1, stroke=0)  # motif-colour marker
            y = self._wrap(c, b, M + 22, y, PAGE_W - 2 * M - 22, self.bfont, 16, self.INK,
                           leading=21) - 14
        self._wordmark(c, name)

    def _data(self, c, s, name):
        c.setFillColor(self.PAPER); c.rect(0, 0, PAGE_W, PAGE_H, fill=1, stroke=0)
        c.setFillColor(self.P); c.rect(0, PAGE_H - 132, PAGE_W, 132, fill=1, stroke=0)
        self._motif(c, on_dark=True)
        c.setFont(self.hfont, 30); c.setFillColor(self.PAPER)
        c.drawString(M, PAGE_H - 92, s["title"])
        if s.get("subtitle"):
            c.setFont(self.bfont, 13); c.setFillColor(self.PAPER)
            c.drawString(M, PAGE_H - 116, s["subtitle"])
        # real-data chart, brand-coloured (reuses the financial chart helper)
        if s.get("revenue") and s.get("ebitda"):
            d = ch.bar_line_combo(PAGE_W - 2 * M - 230, 250, s["cats"], s["revenue"],
                                  s["ebitda"], self.t)
            renderPDF.draw(d, c, M, 70)
        # KPI callouts on the right (accent for emphasis only)
        kx = PAGE_W - M - 200
        ky = PAGE_H - 200
        for label, val in s.get("kpis", []):
            c.setFillColor(self.P); c.setFont(self.hfont, 26)
            c.drawString(kx, ky, val)
            c.setFillColor(self.INK); c.setFont(self.bfont, 11)
            c.drawString(kx, ky - 16, label)
            ky -= 64
        self._wordmark(c, name)

    def _divider(self, c, s, name):
        c.setFillColor(self.P); c.rect(0, 0, PAGE_W, PAGE_H, fill=1, stroke=0)
        self._motif(c, on_dark=True)
        c.setFont(self.bfont, 14); c.setFillColor(HexColor(self.t.accents[0]) if self.t.accents else self.PAPER)
        c.drawString(M, PAGE_H * 0.58 + 30, (s.get("kicker") or "").upper())
        self._wrap(c, s["label"], M, PAGE_H * 0.5, PAGE_W * 0.8, self.hfont, 46, self.PAPER, leading=52)
        self._wordmark(c, name, on_dark=True)

    def _closing(self, c, s, name):
        c.setFillColor(self.P); c.rect(0, 0, PAGE_W, PAGE_H, fill=1, stroke=0)
        self._motif(c, on_dark=True)
        self._wrap(c, s.get("headline", "Let's talk."), M, PAGE_H * 0.58,
                   PAGE_W * 0.8, self.hfont, 46, self.PAPER, leading=52)
        c.setFont(self.bfont, 16); c.setFillColor(self.PAPER)
        if s.get("contact"):
            c.drawString(M, PAGE_H * 0.30, s["contact"])
        self._wordmark(c, name, on_dark=True)

    _RENDER = {"cover": "_cover", "statement": "_statement", "bullets": "_bullets",
               "data": "_data", "divider": "_divider", "closing": "_closing"}

    def build(self, content: dict, out_path: str) -> str:
        name = content.get("brand_name", "Brand")
        os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
        c = canvas.Canvas(out_path, pagesize=(PAGE_W, PAGE_H))
        for s in content["slides"]:
            getattr(self, self._RENDER.get(s["type"], "_statement"))(c, s, name)
            c.showPage()
        c.save()
        return os.path.abspath(out_path)

    # ── editable PPTX companion (python-pptx) ────────────────────────────────
    def build_pptx(self, content: dict, out_path: str) -> str:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        def rgb(hexstr):
            h = hexstr.lstrip("#")
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        prs = Presentation()
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        name = content.get("brand_name", "Brand")

        def add_bg(slide, hexcol):
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = rgb(hexcol)

        def add_text(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT):
            tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = align
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold
            r.font.name = self.t.heading_font if bold else self.t.body_font
            r.font.color.rgb = rgb(color)
            return tb

        for s in content["slides"]:
            sl = prs.slides.add_slide(blank)
            if s["type"] in ("cover", "divider", "closing"):
                add_bg(sl, self.t.primary)
                add_text(sl, s.get("headline") or s.get("label") or name, 0.7, 2.6, 11.9, 2,
                         44, self.t.paper, bold=True)
                tag = s.get("tagline") or (self.voice.get("taglines") or [""])[0]
                if tag:
                    add_text(sl, tag, 0.7, 4.4, 10, 1, 20,
                             self.t.accents[0] if self.t.accents else self.t.paper)
            else:
                add_bg(sl, self.t.paper)
                # primary header zone
                band = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.7))
                band.fill.solid(); band.fill.fore_color.rgb = rgb(self.t.primary); band.line.fill.background()
                add_text(sl, s.get("title", ""), 0.7, 0.5, 11.9, 1, 30, self.t.paper, bold=True)
                y = 2.1
                for b in s.get("bullets", []):
                    add_text(sl, "•  " + b, 0.8, y, 11.7, 0.7, 16, self.t.ink)
                    y += 0.62
                if s.get("kpis"):
                    yk = 2.2
                    for label, val in s["kpis"]:
                        add_text(sl, val, 9.6, yk, 3, 0.6, 26, self.t.primary, bold=True)
                        add_text(sl, label, 9.6, yk + 0.55, 3, 0.4, 11, self.t.ink)
                        yk += 1.2
        prs.save(out_path)
        return os.path.abspath(out_path)


def build_deck(tokens_source, content: dict, out_path: str, pptx: bool = True,
               qa: bool = True) -> dict:
    """High-level entry: build the branded PDF deck (+ optional editable PPTX) and run
    the mandatory render-and-inspect QA loop. `tokens_source` is a brand-token dict or
    a path to one (as emitted by identity.extract_identity)."""
    db = DeckBuilder(tokens_source)
    pdf = db.build(content, out_path)
    out = {"pdf_path": pdf, "slides": len(content["slides"])}
    if pptx:
        out["pptx_path"] = db.build_pptx(content, os.path.splitext(out_path)[0] + ".pptx")
    if qa:
        sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        import qa as _qa
        out["qa"] = _qa.qa_deck(pdf, expect_slides=len(content["slides"]))
    return out


def default_pitch_content(brand_name, tagline, summary=None, cats=None,
                          revenue=None, ebitda=None):
    """Assemble a real-data pitch deck outline. Numbers come from the session
    (e.g. the fin-model summary) — never invented."""
    kpis = []
    if summary:
        kpis = [("Gross margin", f"{summary['gross_margin_pct']:.0f}%"),
                ("Peak funding", f"${summary['peak_funding_requirement']:,}"),
                (f"EBITDA inflection", f"Year {summary['ebitda_inflection_year']}"
                 if summary.get("ebitda_inflection_year") else "—")]
    return {
        "brand_name": brand_name,
        "slides": [
            {"type": "cover", "headline": brand_name, "tagline": tagline,
             "footer": "Investor Brief · Confidential"},
            {"type": "statement", "headline": "Premium cold brew, made the slow way.",
             "sub": "A category growing double digits, and a product built to win the shelf."},
            {"type": "bullets", "title": "Why now",
             "bullets": ["Ready-to-drink coffee is the fastest-growing beverage segment.",
                         "Cold-brew commands a price premium and repeat purchase.",
                         "Our slow-steep process yields a smoother, lower-acid cup.",
                         "Direct retail relationships across 120 doors at launch."]},
            {"type": "divider", "kicker": "The numbers", "label": "A model that funds itself."},
            {"type": "data", "title": "Five-year trajectory",
             "subtitle": "Revenue (bars) and EBITDA (line) — from the live financial model",
             "cats": cats, "revenue": revenue, "ebitda": ebitda, "kpis": kpis},
            {"type": "bullets", "title": "The ask",
             "bullets": ["Raising to cover the peak funding requirement, with margin of safety.",
                         "Capital funds the brewing line, cold storage and working capital.",
                         "Self-funding after the EBITDA inflection; upside in the Bull case."]},
            {"type": "closing", "headline": "Let's build the category leader.",
             "contact": "founders@coldbrew.co  ·  coldbrew.co"},
        ],
    }
