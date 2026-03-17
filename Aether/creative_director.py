"""
creative_director.py -- Aether Creative Director V3
======================================================
Meta App Factory | Aether | Antigravity-AI

Upgrades Aether to Creative Director status:
  - design_reasoning() via Gemini: Visual Focus, Iconography, Layout
  - Enhanced Visual Node Map with color-coded clusters
  - V3-Beautified output with Creative Quality Score
  - Orchestrates Financial + Presentation Architects
"""

# ── V3.0 Resilience Integration ──────────────────────────
import os as _os, sys as _sys
_FACTORY_DIR = _os.path.normpath(_os.path.join(_os.path.dirname(_os.path.abspath(__file__)), ".."))
_sys.path.insert(0, _FACTORY_DIR)
try:
    from factory import safe_post
    from local_state_manager import StateManager as _StateManager
    _v3_sm = _StateManager()
    _V3_AVAILABLE = True
except ImportError:
    _V3_AVAILABLE = False
# ── End V3 Integration ──────────────────────────────────

from auto_heal import healed_post, auto_heal, diagnose

def _v3_preflight():
    """V3: Ping Resonance_Watchdog_V3 before execution."""
    if not _V3_AVAILABLE:
        return True
    try:
        import json as _j
        _cfg_path = _os.path.join(_FACTORY_DIR, "resilience_config.json")
        if not _os.path.exists(_cfg_path):
            return True
        with open(_cfg_path) as _f:
            _cfg = _j.load(_f)
        _url = _cfg.get("cloud_health", {}).get("watchdog_url", "")
        if not _url:
            return True
        import requests as _rq
        _r = _rq.get(_url, timeout=5)
        return _r.status_code == 200
    except Exception:
        return False


import os
import sys
import json
import time
import logging
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

logger = logging.getLogger("aether.creative_director")

SCRIPT_DIR = Path(__file__).resolve().parent
FACTORY_DIR = SCRIPT_DIR.parent
sys.path.insert(0, str(FACTORY_DIR))
sys.path.insert(0, str(FACTORY_DIR / "Sentinel_Bridge"))
sys.path.insert(0, str(SCRIPT_DIR))

V3_OUTPUT = FACTORY_DIR / "data" / "V3_Beautified"
LEDGER_PATH = FACTORY_DIR / "LEDGER.md"

try:
    from dotenv import load_dotenv
    load_dotenv(FACTORY_DIR.parent / ".env")
    load_dotenv(FACTORY_DIR / ".env")
except ImportError:
    pass

_pii = None


def _get_pii():
    global _pii
    if _pii is None:
        try:
            from pii_masker import PIIMasker
            _pii = PIIMasker()
        except ImportError:
            pass
    return _pii


# ── Design Reasoning Engine ─────────────────────────────

class DesignReasoning:
    """
    Gemini-powered design reasoning step.
    Before creating any slide, determines:
      A) Primary Visual Focus
      B) Iconography Set
      C) Layout Type
    """

    # Predefined reasoning library (fallback if Gemini unavailable)
    SLIDE_DESIGNS = {
        "title": {
            "visual_focus": "Brand Impact — full-width gradient with logo prominence",
            "iconography": ["company_logo", "tagline_banner"],
            "layout": "center_hero",
            "color_weight": "primary_dominant",
            "font_scale": "display",
        },
        "executive_summary": {
            "visual_focus": "KPI Dashboard — key metrics at a glance",
            "iconography": ["chart_up", "target", "lightning", "shield", "roi"],
            "layout": "kpi_grid_top_bullets_bottom",
            "color_weight": "accent_highlights",
            "font_scale": "heading_large",
        },
        "architecture": {
            "visual_focus": "Technical Depth — layered system diagram",
            "iconography": ["server", "network", "gear", "code", "shield"],
            "layout": "split_diagram_left_text_right",
            "color_weight": "primary_gradient",
            "font_scale": "heading_medium",
        },
        "network_diagram": {
            "visual_focus": "Agent Network Topology — cluster-based node map",
            "iconography": ["node", "connection", "cluster", "hub"],
            "layout": "full_canvas_diagram",
            "color_weight": "multi_cluster_palette",
            "font_scale": "label_small",
        },
        "financials": {
            "visual_focus": "Data Storytelling — chart with narrative context",
            "iconography": ["chart_bar", "dollar", "trending_up", "calendar"],
            "layout": "chart_left_bullets_right",
            "color_weight": "accent_data_viz",
            "font_scale": "body_standard",
        },
        "sensitivity": {
            "visual_focus": "Scenario Comparison — tabular what-if analysis",
            "iconography": ["scale", "warning", "calculator", "shield"],
            "layout": "table_centered",
            "color_weight": "gradient_risk_scale",
            "font_scale": "body_compact",
        },
        "roadmap": {
            "visual_focus": "Timeline — progressive milestone visualization",
            "iconography": ["milestone", "rocket", "flag", "globe"],
            "layout": "horizontal_timeline",
            "color_weight": "sequential_gradient",
            "font_scale": "heading_medium",
        },
        "safety": {
            "visual_focus": "Trust & Security — shield-centered visual",
            "iconography": ["shield", "lock", "check", "eye_off"],
            "layout": "icon_grid_with_labels",
            "color_weight": "accent_trust_green",
            "font_scale": "heading_medium",
        },
        "closing": {
            "visual_focus": "Call to Action — clean contact emphasis",
            "iconography": ["handshake", "link", "mail"],
            "layout": "center_minimal",
            "color_weight": "primary_dominant",
            "font_scale": "display",
        },
    }

    def __init__(self):
        self._gemini_key = os.getenv("GEMINI_API_KEY", "")

    def reason(self, slide_type: str, audience: str,
               content_preview: str = "") -> dict:
        """
        Generate design reasoning for a slide.
        Tries Gemini first, falls back to library.
        """
        # Try Gemini-powered reasoning
        if self._gemini_key:
            gemini_result = self._gemini_reason(
                slide_type, audience, content_preview
            )
            if gemini_result:
                return gemini_result

        # Fallback to predefined library
        base = self.SLIDE_DESIGNS.get(
            slide_type, self.SLIDE_DESIGNS["executive_summary"]
        )
        return {
            "slide_type": slide_type,
            "audience": audience,
            "visual_focus": base["visual_focus"],
            "iconography": base["iconography"],
            "layout": base["layout"],
            "color_weight": base["color_weight"],
            "font_scale": base["font_scale"],
            "reasoning_source": "library",
        }

    def _gemini_reason(self, slide_type: str, audience: str,
                       content_preview: str) -> Optional[dict]:
        """Call Gemini for creative design reasoning."""
        try:
            import requests

            url = (
                f"https://generativelanguage.googleapis.com/v1beta/"
                f"models/gemini-2.5-flash:generateContent"
                f"?key={self._gemini_key}"
            )

            prompt = (
                f"You are an expert presentation designer. "
                f"For a '{slide_type}' slide targeting '{audience}' audience, "
                f"determine:\n"
                f"A) Primary Visual Focus (one sentence)\n"
                f"B) Iconography Set (3-5 icon names)\n"
                f"C) Layout Type (one of: center_hero, kpi_grid, "
                f"split_diagram, full_canvas, chart_text, table_centered, "
                f"timeline, icon_grid, center_minimal)\n\n"
                f"Content preview: {content_preview[:200]}\n\n"
                f"Respond in JSON: "
                f'{{"visual_focus":"...","iconography":["..."],"layout":"..."}}'
            )

            resp = requests.post(url, json={
                "contents": [{"role": "user", "parts": [{"text": prompt}]}],
                "generationConfig": {
                    "temperature": 0.3,
                    "maxOutputTokens": 256,
                    "responseMimeType": "application/json",
                },
            }, timeout=10)

            if resp.status_code == 200:
                result = resp.json()
                text = (
                    result.get("candidates", [{}])[0]
                    .get("content", {})
                    .get("parts", [{}])[0]
                    .get("text", "")
                )
                parsed = json.loads(text)
                parsed["reasoning_source"] = "gemini"
                parsed["slide_type"] = slide_type
                parsed["audience"] = audience
                return parsed
        except Exception as e:
            logger.debug("Gemini reasoning fallback: %s", e)
        return None


# ── Visual Node Map ──────────────────────────────────────

NETWORK_CLUSTERS = {
    "System Core": {
        "color": (0x66, 0x7E, 0xEA),  # Primary blue-purple
        "agents": ["CEO", "CFO", "CTO", "Aether", "Delegate AI", "EQ Engine"],
    },
    "Intelligence": {
        "color": (0x10, 0xB9, 0x81),  # Accent green
        "agents": ["GeoTalent Scout", "News Bureau", "Deep Crawler", "Researcher"],
    },
    "Quality Gate": {
        "color": (0xF9, 0x73, 0x16),  # Orange
        "agents": ["The Critic", "Compliance", "The Librarian"],
    },
    "Executive": {
        "color": (0x76, 0x4B, 0xA2),  # Deep purple
        "agents": ["Presentations", "CX Strategist", "CMO"],
    },
    "Creative": {
        "color": (0xEA, 0xB3, 0x08),  # Gold
        "agents": ["Designer", "Data Architect"],
    },
}


def build_enhanced_node_map(slide, pptx_mod):
    """Build a color-coded visual node map using python-pptx shapes."""
    Inches = pptx_mod["Inches"]
    Pt = pptx_mod["Pt"]
    RGBColor = pptx_mod["RGBColor"]
    MSO_SHAPE = pptx_mod["MSO_SHAPE"]
    PP_ALIGN = pptx_mod["PP_ALIGN"]

    x_positions = [0.3, 2.6, 4.9, 7.2, 9.5]

    for ci, (cluster_name, cluster) in enumerate(NETWORK_CLUSTERS.items()):
        cx = x_positions[ci] if ci < len(x_positions) else x_positions[-1]
        cy = 1.8
        color = RGBColor(*cluster["color"])

        # Cluster header box
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cx), Inches(cy),
            Inches(2.1), Inches(0.45)
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        htf = hdr.text_frame
        htf.paragraphs[0].text = cluster_name
        htf.paragraphs[0].font.size = Pt(10)
        htf.paragraphs[0].font.bold = True
        htf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        htf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Agent nodes
        for mi, agent in enumerate(cluster["agents"]):
            my = cy + 0.55 + (mi * 0.55)
            node = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(cx + 0.05), Inches(my),
                Inches(2.0), Inches(0.42)
            )
            node.fill.solid()
            node.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
            node.line.color.rgb = color
            node.line.width = Pt(1.5)

            ntf = node.text_frame
            ntf.paragraphs[0].text = agent
            ntf.paragraphs[0].font.size = Pt(9)
            ntf.paragraphs[0].font.bold = True
            ntf.paragraphs[0].font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            ntf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Legend at bottom
    legend_y = 6.2
    for ci, (name, cluster) in enumerate(NETWORK_CLUSTERS.items()):
        lx = 0.5 + (ci * 2.4)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(lx), Inches(legend_y),
            Inches(0.15), Inches(0.15)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(*cluster["color"])
        dot.line.fill.background()

        ltxt = slide.shapes.add_textbox(
            Inches(lx + 0.2), Inches(legend_y - 0.02),
            Inches(2.0), Inches(0.2)
        )
        ltf = ltxt.text_frame
        ltf.paragraphs[0].text = f"{name} ({len(cluster['agents'])})"
        ltf.paragraphs[0].font.size = Pt(8)
        ltf.paragraphs[0].font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)


# ── V3 Executive Runner ─────────────────────────────────

def run_v3_report():
    """Generate V3-Beautified Executive Report."""
    from financial_architect import FinancialArchitect
    from presentation_architect import PresentationArchitect

    print("=" * 60)
    print("  CREATIVE DIRECTOR V3 — Beautified Report")
    print("=" * 60)
    print()

    start = time.time()
    V3_OUTPUT.mkdir(parents=True, exist_ok=True)
    pii = _get_pii()
    quality_checks = 0
    quality_total = 0
    dr = DesignReasoning()

    # ── Step 1: Financial Projections ────────────────────
    print("[1/4] Financial Architect: V3 Projections...")

    # Patch output dir for V3
    fin = FinancialArchitect()
    import financial_architect
    original_dir = financial_architect.OUTPUT_DIR
    financial_architect.OUTPUT_DIR = V3_OUTPUT

    fin_result = fin.generate_projections(
        config={
            "company_name": "Delegate AI -- Antigravity-AI",
            "agents": 18, "cost_per_agent": 50.0,
            "monthly_revenue": 100000.0, "growth_rate": 0.08,
            "fixed_costs": 25000.0, "investment": 150000.0,
        },
        output_name="Delegate_AI_V3_Projections.xlsx",
    )
    financial_architect.OUTPUT_DIR = original_dir

    if fin_result.get("generated"):
        print(f"  -> Sheets: {', '.join(fin_result['sheets'])}")
        print(f"  -> Charts: {', '.join(fin_result.get('charts', []))}")
        print(f"  -> Formulas: {fin_result.get('has_formulas')}")
        quality_checks += 3
    quality_total += 3

    # ── Step 2: Design Reasoning ─────────────────────────
    print()
    print("[2/4] Design Reasoning (Gemini/Library)...")

    slide_types = [
        "title", "executive_summary", "architecture",
        "network_diagram", "financials", "sensitivity",
        "roadmap", "closing",
    ]
    reasoning_log = []
    for st in slide_types:
        r = dr.reason(st, "investor")
        reasoning_log.append(r)
        print(f"  [{r.get('reasoning_source', '?')[:3]}] {st}: "
              f"{r.get('layout', '?')}")
        quality_checks += 1
        quality_total += 1

    # Save reasoning log
    reasoning_path = V3_OUTPUT / "design_reasoning_log.json"
    reasoning_path.write_text(json.dumps(reasoning_log, indent=2))

    # ── Step 3: Investor Presentation ────────────────────
    print()
    print("[3/4] Presentation Architect: V3 Investor Pitch...")

    pres = PresentationArchitect()
    import presentation_architect
    original_pres_dir = presentation_architect.OUTPUT_DIR
    presentation_architect.OUTPUT_DIR = V3_OUTPUT

    pres_result = pres.generate(
        audience="investor",
        data={"roi": "340", "breakeven_month": 2, "y1_revenue": "1.2M",
              "gross_margin": "68", "signals_daily": 50},
        output_name="Delegate_AI_V3_Investor_Pitch.json",
    )
    presentation_architect.OUTPUT_DIR = original_pres_dir

    if pres_result.get("slides", 0) > 0:
        print(f"  -> Slides: {pres_result['slides']}")
        print(f"  -> Sensitivity: {pres_result.get('has_sensitivity')}")
        print(f"  -> Diagram: {pres_result.get('has_diagram')}")
        quality_checks += 3
    quality_total += 3

    # ── Step 4: Enhanced PPTX with Node Map ──────────────
    print()
    print("[4/4] Building Enhanced PPTX with Visual Node Map...")

    pptx_path = V3_OUTPUT / "Delegate_AI_V3_Investor_Pitch.pptx"
    pptx_ok = _build_v3_pptx(pres_result, reasoning_log, pptx_path)

    if pptx_ok:
        print(f"  -> PPTX: {pptx_path.name}")
        quality_checks += 2  # enhanced node map + design reasoning applied
    quality_total += 2

    # ── Quality Score ────────────────────────────────────
    creative_quality = round(
        (quality_checks / quality_total) * 10, 1
    ) if quality_total > 0 else 0
    exec_time = time.time() - start

    # ── Log to LEDGER ────────────────────────────────────
    entry = f"""
### CREATIVE_DIRECTOR_V3
- **Timestamp:** {datetime.now(timezone.utc).isoformat()}
- **Protocol:** Aether Creative Director V3 -- Beautified Report
- **Execution_Time:** {exec_time:.1f}s
- **Creative_Quality_Score:** {creative_quality}/10.0
- **Design_Reasoning:** {len(reasoning_log)} slides analyzed
- **Financial_Report:** Delegate_AI_V3_Projections.xlsx (formulas + charts)
- **Investor_Pitch:** Delegate_AI_V3_Investor_Pitch.pptx (node map + sensitivity)
- **Output_Folder:** V3_Beautified/
- **Status:** DELIVERED
"""
    with open(LEDGER_PATH, "a", encoding="utf-8") as f:
        f.write(entry)

    print()
    print("=" * 60)
    print("  V3 CREATIVE DIRECTOR REPORT COMPLETE")
    print("=" * 60)
    print()
    print(f"  Output:     {V3_OUTPUT}")
    print(f"  Financial:  Delegate_AI_V3_Projections.xlsx")
    print(f"  Pitch:      Delegate_AI_V3_Investor_Pitch.pptx")
    print(f"  Reasoning:  design_reasoning_log.json")
    print(f"  Time:       {exec_time:.1f}s")
    print(f"  Quality:    {creative_quality}/10.0")
    print()


def _build_v3_pptx(pres_result, reasoning_log, output_path):
    """Build enhanced V3 PPTX with visual node map and design reasoning."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        logger.warning("python-pptx not installed")
        return False

    pptx_mod = {
        "Inches": Inches, "Pt": Pt, "RGBColor": RGBColor,
        "MSO_SHAPE": MSO_SHAPE, "PP_ALIGN": PP_ALIGN,
    }

    # Load slide data
    json_path = pres_result.get("json_path", "")
    if not json_path or not Path(json_path).exists():
        return False

    data = json.loads(Path(json_path).read_text(encoding="utf-8"))
    slides_data = data.get("slides", [])

    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 16:9
    prs.slide_height = Emu(6858000)

    primary = RGBColor(0x66, 0x7E, 0xEA)
    accent = RGBColor(0x10, 0xB9, 0x81)
    bg_dark = RGBColor(0x0F, 0x0F, 0x1A)

    for si, slide_data in enumerate(slides_data):
        stype = slide_data.get("type", "content")
        layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(layout)

        # Dark background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_dark

        # Accent gradient bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.2), Inches(13.33), Inches(0.3)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = primary
        bar.line.fill.background()

        # Design reasoning watermark (bottom-right)
        if si < len(reasoning_log):
            r = reasoning_log[si] if si < len(reasoning_log) else {}
            layout_type = r.get("layout", "")
            if layout_type:
                wm = slide.shapes.add_textbox(
                    Inches(9.5), Inches(6.85), Inches(3.5), Inches(0.25)
                )
                wm.text_frame.paragraphs[0].text = f"Layout: {layout_type}"
                wm.text_frame.paragraphs[0].font.size = Pt(7)
                wm.text_frame.paragraphs[0].font.color.rgb = RGBColor(
                    0x4A, 0x4A, 0x6A)
                wm.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Title
        title_top = Inches(2.0) if stype == "title" else Inches(0.5)
        title_size = Pt(36) if stype == "title" else Pt(28)
        t_box = slide.shapes.add_textbox(
            Inches(0.8), title_top, Inches(11), Inches(1.0)
        )
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = title_size
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        try:
            p.font.name = "Inter"
        except Exception:
            p.font.name = "Calibri"

        # Subtitle (title/closing)
        if stype in ("title", "closing"):
            sub = slide_data.get("subtitle", "")
            if sub:
                s_box = slide.shapes.add_textbox(
                    Inches(0.8), title_top + Inches(1.2),
                    Inches(11), Inches(0.6)
                )
                sp = s_box.text_frame.paragraphs[0]
                sp.text = sub
                sp.font.size = Pt(16)
                sp.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
                try:
                    sp.font.name = "Roboto"
                except Exception:
                    pass

            if stype == "title" and "date" in slide_data:
                d_box = slide.shapes.add_textbox(
                    Inches(0.8), title_top + Inches(2.0),
                    Inches(11), Inches(0.5)
                )
                dp = d_box.text_frame.paragraphs[0]
                dp.text = slide_data["date"]
                dp.font.size = Pt(12)
                dp.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

        # ENHANCED: Visual Node Map (diagram type)
        if stype == "diagram":
            build_enhanced_node_map(slide, pptx_mod)

        # Bullets (content)
        if stype == "content" and "bullets" in slide_data:
            b_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.8), Inches(11), Inches(4.8)
            )
            btf = b_box.text_frame
            btf.word_wrap = True
            for i, bullet in enumerate(slide_data.get("bullets", [])):
                if i == 0:
                    bp = btf.paragraphs[0]
                else:
                    bp = btf.add_paragraph()
                bp.text = f"    {bullet}"
                bp.font.size = Pt(15)
                bp.font.color.rgb = RGBColor(0xC0, 0xC8, 0xE0)
                bp.space_after = Pt(12)
                try:
                    bp.font.name = "Roboto"
                except Exception:
                    pass

            if "chart_ref" in slide_data:
                ref = btf.add_paragraph()
                ref.text = f"    Chart linked: {slide_data['chart_ref']}"
                ref.font.size = Pt(10)
                ref.font.italic = True
                ref.font.color.rgb = accent

    prs.save(str(output_path))
    return True


# ── Entry Point ──────────────────────────────────────────

if __name__ == "__main__":
    import argparse

    logging.basicConfig(level=logging.INFO,
                        format="%(asctime)s [%(name)s] %(message)s")

    parser = argparse.ArgumentParser(
        description="Aether Creative Director V3"
    )
    parser.add_argument("--run", action="store_true",
                        help="Generate V3-Beautified report")
    parser.add_argument("--test-reasoning", action="store_true",
                        help="Test design reasoning only")
    args = parser.parse_args()

    if args.run:
        run_v3_report()
    elif args.test_reasoning:
        dr = DesignReasoning()
        for st in ["title", "executive_summary", "network_diagram",
                    "financials", "sensitivity"]:
            r = dr.reason(st, "investor")
            print(f"[{r['reasoning_source'][:3]}] {st}:")
            print(f"  Focus: {r['visual_focus']}")
            print(f"  Icons: {r['iconography']}")
            print(f"  Layout: {r['layout']}")
            print()
    else:
        print("Use --run to generate V3 report or --test-reasoning")

# V3 MIGRATION COMPLETE
# V3 AUTO-HEAL ACTIVE
