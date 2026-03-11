"""
presentation_architect.py -- Aether Presentation Architect V2
================================================================
Meta App Factory | Aether | Antigravity-AI

V2 Executive Polish:
  - High-end PPTX template (gradient bg, styled shapes)
  - Architecture Diagram on Slide 3 (18-agent network)
  - Financial data integration on Slide 5
  - Sensitivity Analysis slide (Investor only)
  - AudienceSwitch for Investor vs Customer differentiation
"""

import os
import sys
import json
import logging
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

logger = logging.getLogger("aether.presentations")

SCRIPT_DIR = Path(__file__).resolve().parent
FACTORY_DIR = SCRIPT_DIR.parent
sys.path.insert(0, str(FACTORY_DIR))
sys.path.insert(0, str(FACTORY_DIR / "Sentinel_Bridge"))

OUTPUT_DIR = FACTORY_DIR / "data" / "V2_Executive_Reports"
CREDS_PATH = FACTORY_DIR / "utils" / "auth" / "google_creds.json"

try:
    from dotenv import load_dotenv
    load_dotenv(FACTORY_DIR.parent / ".env")
    load_dotenv(FACTORY_DIR / ".env")
except ImportError:
    pass

_pii = None


def _get_pii():
    global _pii
    if _pii is None:
        try:
            from pii_masker import PIIMasker
            _pii = PIIMasker()
        except ImportError:
            pass
    return _pii


def get_google_credentials() -> dict:
    return {
        "client_id": os.getenv("GOOGLE_CLIENT_ID", ""),
        "client_secret": os.getenv("GOOGLE_CLIENT_SECRET", ""),
        "project_id": os.getenv("GOOGLE_PROJECT_ID", ""),
        "redirect_uri": os.getenv("GOOGLE_REDIRECT_URI", ""),
        "creds_file": str(CREDS_PATH) if CREDS_PATH.exists() else None,
        "has_credentials": bool(os.getenv("GOOGLE_CLIENT_ID")),
    }


BRAND = {
    "company": "Antigravity-AI",
    "tagline": "Autonomous Intelligence Infrastructure",
    "primary_rgb": (0x66, 0x7E, 0xEA),
    "secondary_rgb": (0x76, 0x4B, 0xA2),
    "accent_rgb": (0x10, 0xB9, 0x81),
    "dark_rgb": (0x0F, 0x0F, 0x1A),
    "light_rgb": (0xF8, 0xFA, 0xFC),
    "text_dark_rgb": (0x1E, 0x29, 0x3B),
    "text_light_rgb": (0x64, 0x74, 0x8B),
}


AUDIENCE_SECTIONS = {
    "investor": {
        "focus": ["Scalability", "Market Intelligence", "ROI", "Architecture"],
        "tone": "professional, data-driven",
        "slides": [
            {"title": "Executive Summary", "func": "_slide_exec_summary"},
            {"title": "Scalable Architecture", "func": "_slide_architecture"},
            {"title": "Agent Network Diagram", "func": "_slide_network_diagram"},
            {"title": "Financial Projections", "func": "_slide_financials"},
            {"title": "Sensitivity Analysis", "func": "_slide_sensitivity"},
            {"title": "Growth Roadmap", "func": "_slide_roadmap"},
        ],
    },
    "customer": {
        "focus": ["User Experience", "Safety & Privacy", "Reliability"],
        "tone": "warm, approachable",
        "slides": [
            {"title": "Welcome", "func": "_slide_welcome"},
            {"title": "Intuitive Dashboard", "func": "_slide_ux"},
            {"title": "Agent Network Diagram", "func": "_slide_network_diagram"},
            {"title": "Your Data is Safe", "func": "_slide_safety"},
            {"title": "Getting Started", "func": "_slide_onboarding"},
        ],
    },
}


class PresentationArchitect:
    """
    V2 Presentation Architect — High-end PPTX with embedded diagrams.
    """

    def __init__(self):
        self._pii = _get_pii()
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    def generate(self, audience: str, data: dict = None,
                 output_name: str = None) -> dict:
        data = data or {}
        audience = audience.lower()
        template = AUDIENCE_SECTIONS.get(audience)
        if not template:
            return {"error": f"Unknown audience: {audience}"}

        # Build slide content
        slides = []

        # Title slide
        slides.append({
            "type": "title",
            "title": f"{BRAND['company']} - {audience.title()} Brief",
            "subtitle": BRAND["tagline"],
            "date": datetime.now().strftime("%B %d, %Y"),
        })

        # Content slides
        for section in template["slides"]:
            func = getattr(self, section["func"], None)
            if func:
                slide = func(data)
                slides.append(slide)

        # Closing
        slides.append({
            "type": "closing",
            "title": "Thank You",
            "subtitle": f"{BRAND['company']} | {BRAND['tagline']}",
        })

        # PII mask
        if self._pii:
            for slide in slides:
                for k in ("title", "subtitle"):
                    if k in slide and isinstance(slide[k], str):
                        slide[k] = self._pii.mask(slide[k])
                if "bullets" in slide:
                    slide["bullets"] = [self._pii.mask(b) for b in slide["bullets"]]

        # Save JSON
        json_name = output_name or f"pitch_{audience}_v2.json"
        json_path = OUTPUT_DIR / json_name
        json_path.write_text(json.dumps({
            "audience": audience,
            "tone": template["tone"],
            "focus": template["focus"],
            "slides": slides,
            "generated_at": datetime.now(timezone.utc).isoformat(),
            "version": "V2",
        }, indent=2))

        # Build PPTX
        pptx_name = json_name.replace(".json", ".pptx")
        pptx_path = OUTPUT_DIR / pptx_name
        pptx_ok = self._build_pptx(slides, audience, pptx_path)

        return {
            "json_path": str(json_path),
            "pptx_path": str(pptx_path) if pptx_ok else None,
            "audience": audience,
            "slides": len(slides),
            "focus": template["focus"],
            "tone": template["tone"],
            "has_sensitivity": audience == "investor",
            "has_diagram": True,
        }

    # ── Slide Content Functions ──────────────────────────

    def _slide_exec_summary(self, data):
        return {
            "type": "content",
            "title": "Executive Summary",
            "bullets": [
                "18-agent autonomous AI network (V7 Router)",
                f"Processing {data.get('signals_daily', 50)}+ actionable signals/day",
                "Self-healing infrastructure with 99.7% uptime target",
                f"Break-even: Month {data.get('breakeven_month', 2)}",
                f"Projected ROI: {data.get('roi', '340')}% Year 1",
            ],
        }

    def _slide_architecture(self, data):
        return {
            "type": "content",
            "title": "Scalable Agent Architecture",
            "bullets": [
                "V7 Agent Skills Router -> 18 specialist agents",
                "Delegate AI Orchestrator for performance-based routing",
                "Leitner spaced-repetition error memory (5-level triage)",
                "DataValidationEngine: Signal Warnings on bad data",
                "News Bureau Chief: auto-triggered context refresh",
            ],
        }

    def _slide_network_diagram(self, data):
        """Architecture diagram slide with 18-agent network layout."""
        agents = [
            ("CEO", "Core"), ("CFO", "Core"), ("CTO", "Core"),
            ("CMO", "Growth"), ("Deep Crawler", "Research"),
            ("Critic", "Quality"), ("Librarian", "Knowledge"),
            ("Compliance", "Safety"), ("Data Architect", "Infra"),
            ("Researcher", "Research"), ("Designer", "Creative"),
            ("Presentations", "Executive"), ("CX Strategist", "Growth"),
            ("Aether", "Core"), ("Delegate AI", "Core"),
            ("EQ Engine", "Core"), ("GeoTalent", "Intelligence"),
            ("News Bureau", "Intelligence"),
        ]
        return {
            "type": "diagram",
            "title": "Agent Network Architecture (V7)",
            "diagram_type": "network",
            "agents": agents,
            "center": "Delegate AI Orchestrator",
            "clusters": {
                "Core": ["CEO", "CFO", "CTO", "Aether", "Delegate AI", "EQ Engine"],
                "Intelligence": ["GeoTalent", "News Bureau", "Deep Crawler", "Researcher"],
                "Quality": ["Critic", "Compliance", "Librarian"],
                "Executive": ["Presentations", "CX Strategist", "CMO"],
                "Creative": ["Designer", "Data Architect"],
            },
        }

    def _slide_financials(self, data):
        return {
            "type": "content",
            "title": "Financial Projections",
            "bullets": [
                f"Break-even: Month {data.get('breakeven_month', 2)}",
                f"Year 1 Revenue: ${data.get('y1_revenue', '1.2M')}",
                f"Gross Margin: {data.get('gross_margin', '68')}%",
                "Cost per agent: $0.02/decision (token-efficient)",
                "See attached XLSX: 12-month P&L + Charts",
            ],
            "chart_ref": "Delegate_AI_2026_Projections_V2.xlsx",
        }

    def _slide_sensitivity(self, data):
        """INVESTOR ONLY: Sensitivity Analysis."""
        return {
            "type": "content",
            "title": "Sensitivity Analysis",
            "subtitle": "What happens when costs change?",
            "bullets": [
                "Base Case: 74.9% margin, Month 2 break-even",
                "Costs +50%: Margin shrinks but remains profitable",
                "Costs Doubled: Still positive annual net profit",
                "Revenue -20%: Margin compressed, breakeven delayed",
                "Full scenario modeling in attached spreadsheet",
            ],
        }

    def _slide_roadmap(self, data):
        return {
            "type": "content",
            "title": "Growth Roadmap",
            "bullets": [
                "Q2 2026: Multi-tenant SaaS deployment",
                "Q3 2026: API marketplace for agent skills",
                "Q4 2026: Enterprise compliance certification",
                "2027: International expansion",
            ],
        }

    def _slide_welcome(self, data):
        return {
            "type": "content",
            "title": "Welcome to Antigravity-AI",
            "bullets": [
                "Your intelligent assistant that adapts to you",
                "18 specialized agents working seamlessly",
                "Beautiful dashboards, simple controls",
                "Your data stays private -- always",
            ],
        }

    def _slide_ux(self, data):
        return {
            "type": "content",
            "title": "Intuitive Dashboard Experience",
            "bullets": [
                "Omni-Dashboard: see all agents at a glance",
                "Real-time stress monitoring adapts tone",
                "One-click context refresh",
                "Dark mode for extended use",
            ],
        }

    def _slide_safety(self, data):
        return {
            "type": "content",
            "title": "Your Data is Safe",
            "bullets": [
                "PIIMasker: auto-redaction on all outputs",
                "FernetVault: encrypted API key storage",
                "No external data leaves unmasked",
                "Compliance-ready audit trails",
            ],
        }

    def _slide_onboarding(self, data):
        return {
            "type": "content",
            "title": "Getting Started",
            "bullets": [
                "1. Launch Meta App Factory dashboard",
                "2. Review agent network health",
                "3. Set notification preferences",
                "4. Let the system work for you",
            ],
        }

    # ── PPTX Builder ─────────────────────────────────────

    def _build_pptx(self, slides, audience, output_path):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.enum.shapes import MSO_SHAPE
        except ImportError:
            logger.warning("python-pptx not installed")
            return False

        prs = Presentation()
        prs.slide_width = Emu(12192000)   # 16:9
        prs.slide_height = Emu(6858000)

        primary = RGBColor(*BRAND["primary_rgb"])
        secondary = RGBColor(*BRAND["secondary_rgb"])
        accent = RGBColor(*BRAND["accent_rgb"])
        dark = RGBColor(*BRAND["text_dark_rgb"])
        light_text = RGBColor(*BRAND["text_light_rgb"])
        bg_dark = RGBColor(*BRAND["dark_rgb"])

        for slide_data in slides:
            stype = slide_data.get("type", "content")
            layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(layout)

            # Dark background
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = bg_dark

            # Accent bar at bottom
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(7.2), Inches(13.33), Inches(0.3)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = primary
            bar.line.fill.background()

            # Title
            title_top = Inches(0.6) if stype != "title" else Inches(2.0)
            title_size = Pt(36) if stype == "title" else Pt(28)
            title_box = slide.shapes.add_textbox(
                Inches(0.8), title_top, Inches(11), Inches(1.2)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.get("title", "")
            p.font.size = title_size
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                p.font.name = "Inter"
            except Exception:
                p.font.name = "Calibri"

            # Subtitle
            if stype in ("title", "closing"):
                sub = slide_data.get("subtitle", "")
                if sub:
                    sub_box = slide.shapes.add_textbox(
                        Inches(0.8),
                        title_top + Inches(1.2),
                        Inches(11), Inches(0.6)
                    )
                    sp = sub_box.text_frame.paragraphs[0]
                    sp.text = sub
                    sp.font.size = Pt(16)
                    sp.font.color.rgb = light_text
                    try:
                        sp.font.name = "Roboto"
                    except Exception:
                        pass

                if stype == "title" and "date" in slide_data:
                    d_box = slide.shapes.add_textbox(
                        Inches(0.8),
                        title_top + Inches(2.0),
                        Inches(11), Inches(0.5)
                    )
                    dp = d_box.text_frame.paragraphs[0]
                    dp.text = slide_data["date"]
                    dp.font.size = Pt(12)
                    dp.font.color.rgb = light_text

            # Network Diagram
            if stype == "diagram":
                clusters = slide_data.get("clusters", {})
                x_start = 0.5
                y_start = 2.0
                col_width = 2.3
                row_height = 0.6

                for ci, (cluster_name, members) in enumerate(clusters.items()):
                    # Cluster header
                    cx = x_start + (ci * col_width)
                    cy = y_start

                    hdr = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(cx), Inches(cy),
                        Inches(2.1), Inches(0.4)
                    )
                    hdr.fill.solid()
                    hdr.fill.fore_color.rgb = primary
                    hdr.line.fill.background()
                    htf = hdr.text_frame
                    htf.paragraphs[0].text = cluster_name
                    htf.paragraphs[0].font.size = Pt(10)
                    htf.paragraphs[0].font.bold = True
                    htf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    htf.paragraphs[0].alignment = PP_ALIGN.CENTER

                    # Agent nodes
                    for mi, member in enumerate(members):
                        my = cy + 0.5 + (mi * row_height)
                        node = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(cx + 0.1), Inches(my),
                            Inches(1.9), Inches(0.4)
                        )
                        node.fill.solid()
                        node.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
                        node.line.color.rgb = primary
                        node.line.width = Pt(1)
                        ntf = node.text_frame
                        ntf.paragraphs[0].text = member
                        ntf.paragraphs[0].font.size = Pt(9)
                        ntf.paragraphs[0].font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
                        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Bullets
            if stype == "content" and "bullets" in slide_data:
                body_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(2.0), Inches(11), Inches(4.5)
                )
                btf = body_box.text_frame
                btf.word_wrap = True
                for i, bullet in enumerate(slide_data.get("bullets", [])):
                    if i == 0:
                        bp = btf.paragraphs[0]
                    else:
                        bp = btf.add_paragraph()
                    bp.text = f"    {bullet}"
                    bp.font.size = Pt(15)
                    bp.font.color.rgb = RGBColor(0xC0, 0xC8, 0xE0)
                    bp.space_after = Pt(12)
                    try:
                        bp.font.name = "Roboto"
                    except Exception:
                        pass

                # Chart reference note
                if "chart_ref" in slide_data:
                    ref = btf.add_paragraph()
                    ref.text = f"    Linked: {slide_data['chart_ref']}"
                    ref.font.size = Pt(10)
                    ref.font.italic = True
                    ref.font.color.rgb = accent

        prs.save(str(output_path))
        logger.info("V2 PPTX generated: %s", output_path.name)
        return True


# ── Entry Point ──────────────────────────────────────────

if __name__ == "__main__":
    import argparse

    logging.basicConfig(level=logging.INFO,
                        format="%(asctime)s [%(name)s] %(message)s")

    parser = argparse.ArgumentParser(
        description="Aether Presentation Architect V2"
    )
    parser.add_argument("--test", action="store_true")
    parser.add_argument("--audience", type=str, default="investor",
                        choices=["investor", "customer"])
    args = parser.parse_args()

    architect = PresentationArchitect()

    if args.test:
        print(f"Presentation Architect V2 -- {args.audience.title()}")
        print("-" * 50)

        result = architect.generate(
            audience=args.audience,
            data={"roi": "340", "breakeven_month": 2, "y1_revenue": "1.2M",
                  "gross_margin": "68", "signals_daily": 50},
        )

        print(f"Audience: {result['audience']}")
        print(f"Slides: {result['slides']}")
        print(f"Tone: {result['tone']}")
        print(f"Focus: {', '.join(result['focus'])}")
        print(f"Has Sensitivity: {result['has_sensitivity']}")
        print(f"Has Diagram: {result['has_diagram']}")
        print(f"JSON: {result['json_path']}")
        if result.get("pptx_path"):
            print(f"PPTX: {result['pptx_path']}")

        print("\nDone!")
    else:
        print("Use --test --audience investor|customer")
